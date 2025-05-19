# 
# PETROS PAPATHEODOROU - 2025
# 

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image
from natsort import natsorted 

print("""
🖼️ PowerPoint Generator from Images 

This script loads JPEG and TIFF images and creates a PowerPoint presentation.
By default, it uses **natural sorting** (like Windows Explorer), for filenames.
    10.jpg  20.jpg  100.jpg

Alternatively, you can override this by running the script with the argument `altsort`:
    `python pptx-from-images.py altsort`
This will sort the filenames lexicographically (alphabetically by characters):
    10.jpg  100.jpg  20.jpg
      
🎨 Background Color:
By default, the slide background is black.
  Add `white` as an argument, to make it white instead.
""")

args = [arg.lower() for arg in sys.argv[1:]]
use_alt_sort = 'altsort' in args
use_white_bg = 'white' in args

use_alt_sort = len(sys.argv) > 1 and sys.argv[1].lower() == 'altsort'

print("Source directory:")
src_dir = input()
print("Generating...")

image_files = [
    f for f in os.listdir(src_dir)
    if f.lower().endswith(('.jpg', '.jpeg', '.tif', '.tiff'))
]

if len(image_files) == 0:
    print("No jpeg or tiff files founded in the provided directory.")
    exit(1)


if use_alt_sort:
    image_files.sort()  # Lexicographic sort
    print("!!! Using alternate (lexicographic) sort")
else:
    image_files = natsorted(image_files)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide_w = prs.slide_width
slide_h = prs.slide_height

for image_file in image_files:
    img_path = os.path.join(src_dir, image_file)

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    img_ratio = img_w / img_h
    slide_ratio = slide_w / slide_h

    if img_ratio > slide_ratio:
        # Image is wider → match width
        new_width = slide_w
        new_height = int(new_width / img_ratio)
    else:
        # Image is taller → match height
        new_height = slide_h
        new_width = int(new_height * img_ratio)

    # Center image on slide
    left = (slide_w - new_width) / 2
    top = (slide_h - new_height) / 2

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Set slide background to black
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255) if use_white_bg else RGBColor(0, 0, 0)
    slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

prs.save("output_presentation.pptx")
print("output_presentation.pptx saved.")
